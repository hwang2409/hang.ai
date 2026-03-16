import io
import re
from html.parser import HTMLParser
from urllib.parse import urlparse


async def extract_pdf_text(file_bytes: bytes) -> str:
    """Extract text from PDF bytes."""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
    return "\n\n".join(pages)


async def extract_pptx_text(file_bytes: bytes) -> str:
    """Extract text from PPTX bytes."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


async def extract_youtube_transcript(url: str) -> tuple[str, str]:
    """Extract transcript from YouTube URL. Returns (transcript_text, video_title)."""
    from youtube_transcript_api import YouTubeTranscriptApi

    # Extract video ID
    video_id = None
    patterns = [
        r'(?:v=|/v/|youtu\.be/)([a-zA-Z0-9_-]{11})',
        r'(?:embed/)([a-zA-Z0-9_-]{11})',
    ]
    for pat in patterns:
        match = re.search(pat, url)
        if match:
            video_id = match.group(1)
            break

    if not video_id:
        raise ValueError("Could not extract video ID from URL")

    ytt = YouTubeTranscriptApi()

    # Try fetching English transcript directly first
    entries = None
    try:
        entries = ytt.fetch(video_id, languages=['en'])
    except Exception:
        # Fall back to listing and finding any available transcript
        try:
            transcript_list = ytt.list(video_id)
            try:
                transcript = transcript_list.find_manually_created_transcript(['en'])
            except Exception:
                try:
                    transcript = transcript_list.find_generated_transcript(['en'])
                except Exception:
                    transcript = None
                    for t in transcript_list:
                        transcript = t
                        break
            if transcript:
                entries = transcript.fetch()
        except Exception:
            pass

    if not entries:
        raise ValueError("No transcript available for this video")
    # Group into ~60 second chunks for readability
    chunks = []
    current_chunk = []
    chunk_start = 0

    for entry in entries:
        if entry.start - chunk_start > 60 and current_chunk:
            mins = int(chunk_start // 60)
            secs = int(chunk_start % 60)
            timestamp = f"[{mins:02d}:{secs:02d}]"
            chunks.append(f"{timestamp} " + " ".join(current_chunk))
            current_chunk = []
            chunk_start = entry.start
        current_chunk.append(entry.text)

    if current_chunk:
        mins = int(chunk_start // 60)
        secs = int(chunk_start % 60)
        timestamp = f"[{mins:02d}:{secs:02d}]"
        chunks.append(f"{timestamp} " + " ".join(current_chunk))

    text = "\n\n".join(chunks)
    title = f"YouTube Video {video_id}"

    # Try to get title via a simple fetch
    try:
        import httpx
        resp = await httpx.AsyncClient(timeout=10).get(url)
        match = re.search(r'<title>([^<]+)</title>', resp.text)
        if match:
            raw_title = match.group(1).replace(" - YouTube", "").strip()
            if raw_title:
                title = raw_title
    except Exception:
        pass

    return text, title


def detect_url_type(url: str) -> str:
    """Detect URL type: 'youtube', 'arxiv', or 'webpage'."""
    host = urlparse(url).hostname or ""
    host = host.lower()
    if any(d in host for d in ("youtube.com", "youtu.be")):
        return "youtube"
    if "arxiv.org" in host:
        return "arxiv"
    return "webpage"


def _extract_arxiv_id(url: str) -> str:
    """Extract arXiv paper ID from URL."""
    # Handles abs/2301.12345, pdf/2301.12345, html/2301.12345v2
    match = re.search(r'(?:abs|pdf|html)/(\d{4}\.\d{4,5}(?:v\d+)?)', url)
    if match:
        return match.group(1)
    # Handles arxiv.org/2301.12345
    match = re.search(r'arxiv\.org/(\d{4}\.\d{4,5}(?:v\d+)?)', url)
    if match:
        return match.group(1)
    raise ValueError(f"Could not extract arXiv ID from URL: {url}")


async def extract_arxiv_content(url: str) -> tuple[str, str, dict]:
    """Extract content from arXiv abs page via meta tags (avoids rate-limited API)."""
    import httpx

    arxiv_id = _extract_arxiv_id(url)
    abs_url = f"https://arxiv.org/abs/{arxiv_id}"

    async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
        resp = await client.get(abs_url, headers={"User-Agent": "Mozilla/5.0 (compatible; HangBot/1.0)"})
        resp.raise_for_status()

    html = resp.text

    # Extract from citation_* meta tags
    def meta(name):
        m = re.search(rf'<meta\s+name="{name}"\s+content="([^"]*)"', html)
        return m.group(1).strip() if m else ""

    def meta_all(name):
        return [m.group(1).strip() for m in re.finditer(rf'<meta\s+name="{name}"\s+content="([^"]*)"', html)]

    title = meta("citation_title") or f"arXiv:{arxiv_id}"
    authors = meta_all("citation_author")
    abstract = meta("citation_abstract")
    # og:description as fallback for abstract
    if not abstract:
        m = re.search(r'<meta\s+property="og:description"\s+content="([^"]*)"', html)
        abstract = m.group(1).strip() if m else ""

    pdf_url = f"https://arxiv.org/pdf/{arxiv_id}"

    text = f"# {title}\n\nAuthors: {', '.join(authors)}\n\n## Abstract\n\n{abstract}"

    metadata = {
        "title": title,
        "authors": authors,
        "pdf_url": pdf_url,
        "arxiv_id": arxiv_id,
        "domain": "arxiv.org",
        "thumbnail_url": None,
        "description": abstract[:300] if abstract else None,
    }

    return text, title, metadata


class _SimpleHTMLTextExtractor(HTMLParser):
    """Minimal HTML parser to extract readable text."""

    SKIP_TAGS = {"script", "style", "noscript", "svg", "nav", "footer", "header"}

    def __init__(self):
        super().__init__()
        self.result = []
        self._skip_depth = 0
        self.title = ""
        self._in_title = False
        self.og_description = ""
        self.og_image = ""

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        if tag in self.SKIP_TAGS:
            self._skip_depth += 1
        if tag == "title":
            self._in_title = True
        if tag == "meta":
            prop = attrs_dict.get("property", "")
            name = attrs_dict.get("name", "")
            content = attrs_dict.get("content", "")
            if prop == "og:description" or name == "description":
                if not self.og_description:
                    self.og_description = content
            if prop == "og:image":
                self.og_image = content

    def handle_endtag(self, tag):
        if tag in self.SKIP_TAGS:
            self._skip_depth = max(0, self._skip_depth - 1)
        if tag == "title":
            self._in_title = False
        if tag in ("p", "div", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr"):
            self.result.append("\n")

    def handle_data(self, data):
        if self._in_title:
            self.title += data
        if self._skip_depth == 0:
            self.result.append(data)

    def get_text(self) -> str:
        raw = "".join(self.result)
        # Collapse whitespace within lines, preserve paragraph breaks
        lines = raw.split("\n")
        cleaned = []
        for line in lines:
            line = " ".join(line.split())
            if line:
                cleaned.append(line)
        return "\n\n".join(cleaned)


async def extract_webpage_content(url: str) -> tuple[str, str, dict]:
    """Extract content from a webpage URL. Returns (text, title, metadata_dict)."""
    import httpx

    async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
        resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 (compatible; HangBot/1.0)"})
        resp.raise_for_status()

    parser = _SimpleHTMLTextExtractor()
    parser.feed(resp.text)
    text = parser.get_text()
    title = parser.title.strip() or urlparse(url).hostname or url

    # Truncate very long pages
    if len(text) > 50000:
        text = text[:50000] + "\n\n[Content truncated]"

    domain = urlparse(url).hostname or ""
    metadata = {
        "title": title,
        "domain": domain,
        "description": parser.og_description[:300] if parser.og_description else (text[:300] if text else None),
        "thumbnail_url": parser.og_image or None,
    }

    return text, title, metadata
